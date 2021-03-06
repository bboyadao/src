import os
import sys
import csv
import time
import json
import base64
import zipfile
import pathlib
import requests
import xmltodict
import subprocess
from os import listdir
# import comtypes.client
from io import BytesIO
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfFileMerger
from os.path import isfile, join
from distutils.dir_util import copy_tree
from shutil import (copy2, copytree, rmtree, move)
API_KEY = "a9452d93d8b18b78fec035be138daebc"
PASSWORD = "ee2fb8b0b04e52c48e7ee6c61586c176"
SHOP = "taolaadao"
url = f"https://{API_KEY}:{PASSWORD}@{SHOP}.myshopify.com/admin/api/2019-04/products.json"
_encode = base64.b64encode(bytes(url, 'utf-8')).decode('ascii')
headers = {"Authorization": f"Basic {_encode}"}
owl_id = "3387976189ccd7b"
owl_key = "c9c0214dc44a8efd567a"
owlurl = f'https://{owl_id}:{owl_key}@upload.sendowl.com/api/v1/products.xml'
from operator import itemgetter

wdFormatPDF = 17


def zipdir(path, ziph):
    length = len(path)

    # ziph is zipfile handle
    for root, dirs, files in os.walk(path):
        folder = root[length:]  # path without "parent"
        for file in files:
            ziph.write(os.path.join(root, file), os.path.join(folder, file))


BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

mypath = os.path.join(BASE_DIR)
src_path = os.path.join(BASE_DIR, "src")

templates = os.path.join(BASE_DIR, src_path, "Templates")
shopify_copy_course = os.path.join(BASE_DIR, templates, "Shopify",
                                   "shopify_copy_course.txt")
shopify_copy_book = os.path.join(BASE_DIR, templates, "Shopify",
                                 "shopify_copy_book.txt")
SEO_path = os.path.join(BASE_DIR, templates, "Shopify",
                        "shopify_SEO_company.txt")


onlyfiles = [f for f in listdir(templates) if isfile(join(templates, f))]

onlyfolder = [f for f in listdir(templates) if not isfile(join(templates, f))]


# def publish_product(p_id):
#     url = f"https://{API_KEY}:{PASSWORD}@{SHOP}.myshopify.com/admin/api/2019-04/products{p_id}.json"
#     data = {
#         "product": {
#             "id": p_id,
#             "published": True
#         }
#     }
#
#     r = requests.post(url, json=data, headers=headers)

def update_requires_shiping(v_id):
    url = f"https://{API_KEY}:{PASSWORD}@{SHOP}.myshopify.com/admin/api/2019-04/variants/{v_id}.json"
    data = {
        "variant": {
            "id": v_id,
            "requires_shipping": False
        }
    }
    r = requests.put(url, json=data, headers=headers)
    if r.status_code != 200:
        print(r.content)
        return False
    return True


def add_product(cp_name, position, industry, price, title, path):
    url = f"https://{API_KEY}:{PASSWORD}@{SHOP}.myshopify.com/admin/api/2019-04/products.json"

    with open(path, encoding="utf-8") as r:
        f = r.read()
    if "[Company]" or "[company]" or "[Position]" or "[position]" or "[Industry]" or "[industry]" in f:
        des1 = f.replace("[company]", cp_name)
        des2 = des1.replace("[position]", position)
        des3 = des2.replace("[Company]", cp_name)
        des4 = des3.replace("[Position]", position)
        des5 = des4.replace("[Industry]", position)
        des6 = des5.replace("[industry]", position)

    des = des6
    with open(SEO_path, encoding="utf-8") as r:
        f = r.read()
    if "[Company]" or "[company]" or "[Position]" or "[position]"in f:
        seo1 = f.replace("[company]", cp_name)
        seo2 = seo1.replace("[position]", position)
        seo3 = seo2.replace("[Company]", cp_name)
        seo4 = seo3.replace("[Position]", position)
    seo = seo4

    product = {
        "product": {
            "metafields_global_title_tag": seo,
            "metafields_global_description_tag": seo,
            "title": title,
            "published": "true",

            "body_html": des,
            "vendor": "Coursetake",
            "product_type": "Digital",
            "variants": [{"price": price}]
            # "tags": seo2
        }
    }
    r = requests.post(url, headers=headers, json=product)
    return r.json()


def upload_image(product_id, img_path):
    url = f"https://{API_KEY}:{PASSWORD}@{SHOP}.myshopify.com/admin/api/2019-04/products/{product_id}/images.json"
    with open(img_path, "rb") as f:
        img_str = f.read()

    img_str = base64.b64encode(img_str)

    data = {
        "image": {
            "attachment": str(img_str, 'utf-8'),
            "filename": os.path.basename(img_path)
        }
    }
    r = requests.post(url, json=data, headers=headers)
    return r.json()


def send_owl(title, price, v_id, zip_path, pdf_stamping):

    headers = {
        "Content-type": "multipart/form-data,application/json",
        "Accept": "application/json"
    }
    files = {
        'product[name]': (None, title),
        'product[product_type]': (None, 'digital'),
        'product[price]': (None, price),
        'product[shopify_variant_id]': (None, v_id),
        'product[pdf_stamping]': (None, pdf_stamping),
        'product[attachment]': (os.path.basename(zip_path), open(zip_path, 'rb')),
    }
    r = requests.post(owlurl, files=files,)

    if r.status_code != 201:
        print(f"there are error {r.content}")
        return None
    return json.dumps(xmltodict.parse(r.text))


def check_collection(name):
    url = f"https://{API_KEY}:{PASSWORD}@{SHOP}.myshopify.com/admin/api/2019-04/custom_collections.json"
    r = requests.get(url)
    data = r.json()
    for i in data['custom_collections']:
        if i['title'] == name:
            return (True, i['id'])
    return (False, None)


def create_collection(cp_name, logo):

    url = f"https://{API_KEY}:{PASSWORD}@{SHOP}.myshopify.com/admin/api/2019-04/custom_collections.json"
    collect = {
        "custom_collection": {
            "title": cp_name,
            "body": f"Courses and Study Guides to help you ace your upcoming interview at {cp_name}",
            "image": {
                "src": logo,
            },

        }}

    r = requests.post(url, json=collect)
    if r.status_code != 201:
        return False
    return True


def add_product_to_collection(p_id, c_id):
    url = f"https://{API_KEY}:{PASSWORD}@{SHOP}.myshopify.com/admin/api/2019-04/collects.json"
    data = {"collect": {
        "product_id": p_id,
        "collection_id": c_id
    }
    }
    r = requests.post(url, json=data)
    if r.status_code != 201:
        return False
    return True


def update_toc(docx_file):
    word = comtypes.client.CreateObject("Word.Application")
    doc = word.Documents.Open(docx_file)
    time.sleep(5)
    toc_count = doc.TablesOfContents.Count

    if toc_count > 0:
        toc = doc.TablesOfContents(1)

        toc.UpdatePageNumbers
        toc.Update
        time.sleep(10)
        print('TOC should have been updated.')
    else:
        print('TOC has not been updated...')

    doc.Save()
    doc.SaveAs(docx_file, FileFormat=16)
    doc.Close(SaveChanges=True)
    word.Quit()


def merged_by_macro(clone, merged_name):
    active_dir = os.path.dirname(clone)

    _files = [f for f in listdir(active_dir) if isfile(join(active_dir, f))]
    dud = ["co", "~$"]
    test = [[x, int(x[0:2].replace("-", "").replace(" ", ""))]
            for x in _files if x[0:2] not in dud]
    test = sorted(test, key=itemgetter(1))
    files = [x[0] for x in test]

    active_files = ', '.join('""{0}""'.format(w) for w in files)

    macro = '''
Sub NewDocWithCode()
    Dim doc As Document
    Set doc = ActiveDocument
    doc.VBProject.VBComponents("ThisDocument").CodeModule.AddFromString _
    "Sub Mergedocuments()" & vbLf & _
        "Application.ScreenUpdating = False" & vbLf & _
        "MyPath = ActiveDocument.Path" & vbLf & _
        "Dim myHeadings" & vbLf & _
        "ActiveDocument.Characters.Last.Select" & vbLf & _
        "Selection.Collapse" & vbLf & _
        "myHeadings = Array(''' + active_files + ''')" & vbLf & _
        "For Each Heading In myHeadings" & vbLf & _
            "b = MyPath & ""\\"" & Heading" & vbLf & _
            "Set wb = Documents.Open(b)" & vbLf & _
            "Selection.WholeStory" & vbLf & _
            "Selection.Copy" & vbLf & _
            "Windows(1).Activate" & vbLf & _
            "Selection.EndKey Unit:=wdLine" & vbLf & _
            "Selection.TypeParagraph" & vbLf & _
            "Selection.Paste" & vbLf & _
            "wb.Close False" & vbLf & _
            "Next Heading" & vbLf & _
        "End Sub"
End Sub

    '''

    macr2 = """
    Sub Macro1()
    Selection.WholeStory
    Selection.Delete Unit:=wdCharacter, Count:=1
    End Sub
"""

    update_toc_macro = '''
Sub Update_toc()
    Dim t As TableOfContents
    For Each t In ActiveDocument.TablesOfContents
        t.Update
    Next t
    ActiveDocument.Fields.Update
End Sub
     '''
    word = comtypes.client.CreateObject("Word.Application")
    doc = word.Documents.Open(clone)
    wordModule = doc.VBProject.VBComponents.Add(1)
    wordModule.CodeModule.AddFromString(macr2)
    word.Application.Run("Macro1")
    doc.Save()
    doc.SaveAs(clone, FileFormat=16)
    doc.Close()
    word.Quit()

    word = comtypes.client.CreateObject("Word.Application")

    doc = word.Documents.Open(clone)
    time.sleep(5)
    wordModule = doc.VBProject.VBComponents.Add(1)
    wordModule.CodeModule.AddFromString(macro)

    word.Application.Run("NewDocWithCode")
    time.sleep(1)
    word.Application.Run("Mergedocuments")
    time.sleep(10)
    doc.Save()
    doc.SaveAs(merged_name, FileFormat=16)
    doc.Close()
    word.Quit()

    word_toc = comtypes.client.CreateObject("Word.Application")
    doc_toc = word_toc.Documents.Open(merged_name)
    print("===== Updating Table of Contents")
    wordModule = doc_toc.VBProject.VBComponents.Add(1)
    wordModule.CodeModule.AddFromString(update_toc_macro)
    word_toc.Application.Run("Update_toc")
    doc_toc.SaveAs(merged_name, FileFormat=16)
    doc_toc.Close()
    word_toc.Quit()

    _path = os.path.dirname(clone)
    _files = [f for f in listdir(_path) if isfile(join(_path, f))]
    for i in _files:
        i = os.path.join(_path, i)
    return _files


def replace_word(path, cp_name, position, industry, logo):
    doc = Document(path)
    header = doc.sections[0].header

    for i, j in enumerate(header.paragraphs):

        if '[Company]' in j.text:
            header.paragraphs[i].text = j.text.replace(
                '[Company]', cp_name)
        if '[Position]' in j.text:
            header.paragraphs[i].text = j.text.replace(
                '[Position]', position)
    for g, k in enumerate(doc.paragraphs):
        if '[Company]'in k.text:
            doc.paragraphs[g].text = k.text.replace(
                '[Company]', cp_name)

        if'[Position]' in k.text:
            doc.paragraphs[g].text = k.text.replace(
                '[Position]', position)

        if'[Industry]' in k.text:
            doc.paragraphs[g].text = k.text.replace(
                '[Industry]', industry)

        if'[Company Image]' in k.text:
            doc.paragraphs[g].text = ""
            try:
                response = requests.get(logo)
            except:
                response = ""
            binary_img = BytesIO(response.content)
            pic = doc.paragraphs[g]
            run = pic.add_run()
            run.add_picture(binary_img)
    doc.save(path)
    return doc


def get_form(path):

    with open(os.path.join(src_path, path)) as f:
        b = csv.reader(f)
        return list(b)


def replace_and_convert_pptx_to_pdf(src, dst, data):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    deck = powerpoint.Presentations.Open(src)
    wordModule = deck.VBProject.VBComponents.Add(1)
    replaced = ', '.join('"{0}"'.format(w) for w in data)
    macro = '''
    Sub Multi_FindReplace()
    Dim sld As Slide
    Dim shp As Shape
    Dim ShpTxt As TextRange
    Dim TmpTxt As TextRange
    Dim FindList As Variant
    Dim ReplaceList As Variant
    Dim x As Long
    Dim t As String
    Dim y As String

    FindList = Array("[Company]", "[Position]", "[Industry]")
    ReplaceList = Array(''' + replaced + ''')

    For Each sld In ActivePresentation.Slides

       For x = LBound(FindList) To UBound(FindList)

        t = FindList(x)
        y = ReplaceList(x)

        Call findAndReplaceText(sld, t, y)

       Next x

    Next sld

    End Sub

    Sub findAndReplaceText(sld As PowerPoint.Slide, findText As String, replaceText As String)
    Dim shp As PowerPoint.Shape
    Dim textLoc As PowerPoint.TextRange
    For Each shp In sld.Shapes
        If shp.HasTextFrame Then
            If shp.TextFrame.HasText Then

               Set textLoc = shp.TextFrame.TextRange.Find(findText)
               Do While Not (textLoc Is Nothing)
                If Not (textLoc Is Nothing) Then
                 textLoc.Text = replaceText

                End If
                Set textLoc = shp.TextFrame.TextRange.Find(findText)
               Loop
            End If
        End If
    Next shp
    End Sub

    '''
    wordModule.CodeModule.AddFromString(macro)
    deck.Application.Run("Multi_FindReplace", "")

    deck.SaveAs(dst, 32)
    deck.Close()
    powerpoint.Quit()
    return True


def convert_to_pdf(src, dst):

    dst = dst.replace("docx", "pdf")
    word = comtypes.client.CreateObject('Word.Application')
    word.Visible = False
    doc = word.Documents.Open(src)
    doc.SaveAs(dst, FileFormat=wdFormatPDF)
    doc.Close()
    word.Quit()
    return True


def replace_ppxt(path, new_path, cp_name, position, industry):
    ppt = Presentation(path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                a = shape.text = shape.text.replace('[Company]', cp_name)
                b = shape.text = a.replace('[Position]', position)
                c = shape.text = b.replace('[Industry]', industry)
    ppt.save(new_path)
    return True


def create_sys_temp_dir(files, sys_temp_dir, cp_name, position, industry, logo):
    list_file = []
    for i, j in enumerate(files):

        a = f"{i}-{os.path.basename(j)}"
        clone = f"copy_template.docx"
        if i == 0:

            copy2(j, os.path.join(sys_temp_dir, a))
            copy2(j, os.path.join(sys_temp_dir, clone))
            replace_word(os.path.join(sys_temp_dir, clone),
                         cp_name, position, industry, logo)

        else:
            copy2(j, os.path.join(sys_temp_dir, a))
        list_file.append(os.path.join(sys_temp_dir, a))

    return (sys_temp_dir, list_file)


if __name__ == '__main__':

    file_input = [os.path.join(templates, i) for i in onlyfiles]

    form = get_form("Industry_Company_Position.csv")

    for data in form[1:-1]:

        cp_name = data[0].strip()
        logo = data[1].strip()
        industry = data[3].strip()
        position = data[4].strip()
        questionsandanswers = data[-2].strip()
        interviewprocess = data[-1].strip()
        price = data[5].strip()
        price_course = data[6].strip()

        parent_dir = os.path.join(
            BASE_DIR, "src", "Output", cp_name, position)
        temp_dir = os.path.join(parent_dir, "Temp Files")

        pathlib.Path(temp_dir).mkdir(parents=True, exist_ok=True)

        pathlib.Path(os.path.join(parent_dir, "Study Guide")
                     ).mkdir(parents=True, exist_ok=True)

        pathlib.Path(os.path.join(temp_dir, "sys_temp_dir")
                     ).mkdir(parents=True, exist_ok=True)

        pathlib.Path(os.path.join(parent_dir, "Course")
                     ).mkdir(parents=True, exist_ok=True)

        for i in file_input:
            if "$" in i:
                continue
            copy2(i, temp_dir)
        #
        new_dir = os.path.join(BASE_DIR, "src", "Output", cp_name, position,
                               "Temp Files")
        #
        for i in onlyfolder[0:-1]:
            copy_tree(os.path.join(templates, i), os.path.join(temp_dir, i))
        word = onlyfolder[-1]
        _cp_detail = os.path.join(templates, word, "Company Details")
        _industry = os.path.join(templates, word, "Industry Details")
        _interview_process = os.path.join(templates, word, "Interview Process")
        _jd = os.path.join(templates, word, "Job Description")
        _qa = os.path.join(templates, word, "List of Questions and Answers")
        _img = os.path.join(templates, "Images", "Course")

        cp_detail_word = [f for f in listdir(
            _cp_detail) if not isfile(join(templates, f))]
        industry_word = [f for f in listdir(
            _industry) if not isfile(join(templates, f))]
        interview_process_word = [f for f in listdir(
            _interview_process) if not isfile(join(templates, f))]

        jd_word = [f for f in listdir(
            _jd) if not isfile(join(templates, f))]
        qa_word = [f for f in listdir(
            _qa) if not isfile(join(templates, f))]
        img = [f for f in listdir(
            _img) if not isfile(join(templates, f))]
        print("===== Copying Template Files for " + cp_name)

        for i in cp_detail_word:
            if cp_name.lower() in i.lower():
                copy2(os.path.join(_cp_detail, i), temp_dir)

        for i in industry_word:
            if industry.lower() in i.lower():
                copy2(os.path.join(_industry, i), temp_dir)

        for i in interview_process_word:
            if interviewprocess.lower() in i.lower():
                copy2(os.path.join(_interview_process, i), temp_dir)

        for i in jd_word:

            if f"{cp_name} {position}.docx".lower() == i.lower():
                copy2(os.path.join(_jd, i), temp_dir)

        for i in qa_word:
            if f"{questionsandanswers}.docx".lower() == i.lower():
                copy2(os.path.join(_qa, i), temp_dir)

        for i in img:
            if f"{cp_name} {position}.jpg".lower() == i.lower():
                copy2(os.path.join(_img, i), os.path.join(
                    temp_dir, "Images"))

        now_only_files = [f for f in listdir(
            new_dir) if isfile(join(new_dir, f))]

        only_files_docx = [f for f in now_only_files if "docx" in f]

        only_files_docx = [f for f in only_files_docx if f[:2] != "~$"]

        study_file = [f for f in only_files_docx if "Study" in f]
        workbook_file = [f for f in only_files_docx if "Workbook" in f]
        workbook_file.sort()

        study_list = []
        workbook_list = []
        study_file.sort()

        study_file.insert(1, f"{interviewprocess}.docx")
        study_file.insert(3, f"{industry}.docx")
        study_file.insert(5, f"{cp_name}.docx")
        study_file.insert(7, f"{cp_name} {position}.docx")
        study_file.insert(9, f"{questionsandanswers}.docx")
        merged_study = os.path.join(
            new_dir, f'Study Guide–{cp_name} {position} Interview preparation.docx')
        merged_workbook = os.path.join(
            new_dir, f'Workbook–{cp_name} {position} Interview preparation.docx')
        merged_study_pdf = os.path.join(
            new_dir, f'Study Guide–{cp_name} {position} Interview preparation.pdf')
        merged_workbook_pdf = os.path.join(
            new_dir, f'Workbook–{cp_name} {position} Interview preparation.pdf')

        for i in study_file:
            study_list.append(os.path.join(new_dir, i))

        pdf_study_list = [i.replace("docx", "pdf") for i in study_list]

        (list_temp_dir, list_file) = create_sys_temp_dir(
            study_list, os.path.join(temp_dir, "sys_temp_dir"), cp_name, position, industry, logo)
        _stu = os.path.join(
            temp_dir, f"Study Guide–{cp_name} {position} Interview preparation.docx")

        print("===== Replacing words...")
        # for i in list_file:
        #     replace_word(i, cp_name, position, industry, logo)
        print(
            f"===== Merging Templates into: Study Guide–{cp_name} {position} Interview preparation.docx")
        # list_files_in_temp = merged_by_macro(os.path.join(
        # temp_dir, "sys_temp_dir", "copy_template.docx"), merged_study)

        print(
            f"===== Creating PDF: Study Guide–{cp_name} {position} Interview preparation.pdf")

        # convert_to_pdf(merged_study, merged_study_pdf)

        pathlib.Path(os.path.join(temp_dir, "sys_temp_dir")
                     ).mkdir(parents=True, exist_ok=True)

        # copy2(merged_study_pdf, os.path.join(
        #     parent_dir, "Course", f'Study Guide–{cp_name} {position} Interview preparation.pdf'))
        #
        # copy2(merged_study_pdf, os.path.join(
        #     parent_dir, "Study Guide", f'Study Guide–{cp_name} {position} Interview preparation.pdf'))

        workbook_file.insert(1, f"{interviewprocess}.docx")
        workbook_file.insert(3, f"{industry}.docx")
        workbook_file.insert(5, f"{cp_name}.docx")
        workbook_file.insert(7, f"{cp_name} {position}.docx")
        workbook_file.insert(9, f"{questionsandanswers}.docx")
        for i in workbook_file:
            workbook_list.append(os.path.join(new_dir, i))
        _work = os.path.join(
            temp_dir, f"Workbook–{cp_name} {position} Interview preparation.docx")

        print("===== Replacing words...")
        (list_temp_dir, list_file) = create_sys_temp_dir(
            workbook_list, os.path.join(temp_dir, "sys_temp_dir"), cp_name, position, industry, logo)
        # for i in list_file:
        #     replace_word(i, cp_name, position, industry, logo)
        print(
            f"===== Merging Templates into: Workbook–{cp_name} {position} Interview preparation.docx")
        # list_files_in_temp = merged_by_macro(os.path.join(
        # temp_dir, "sys_temp_dir", "copy_template.docx"), merged_workbook)

        print(
            f"===== Creating PDF: Workbook–{cp_name} {position} Interview preparation.pdf")
        # convert_to_pdf(merged_workbook, merged_workbook_pdf)

        # copy2(merged_workbook_pdf, os.path.join(
        #     parent_dir, "Course", f'Workbook–{cp_name} {position} Interview preparation.pdf'))
        for i in onlyfiles:
            if "pdf" in i:
                if f"Interview preparation" in i:
                    continue
                else:
                    os.remove(os.path.join(new_dir, i))

        print("===== Converting Powerpoints to pdf")
        pptx_file = os.path.join(
            new_dir, "Course", "Slides - Coursetake Interview Preparation.pptx")
        pptx_to = os.path.join(
            new_dir, "Course", f"Slides – {cp_name} {position} Interview preparation.pptx")
        # os.rename(pptx_file, pptx_to)

        # replace_ppxt(pptx_to, cp_name, position, industry)
        print(pptx_to)
        # replace_and_convert_pptx_to_pdf(pptx_to,
        #                                 os.path.join(
        #                                     parent_dir,
        #                                     "Course",
        #                                     f"Slides – {cp_name} {position} Interview preparation.pdf"
        #                                 ),
        #                                 [cp_name, position, industry]
        #                                 )

        src_course = os.path.join(new_dir, "Course")
        src_file_course = [f for f in listdir(
            src_course) if isfile(join(src_course, f))]
        for i in src_file_course:
            copy2(os.path.join(src_course, i),
                  os.path.join(parent_dir, "Course"))
        os.chdir(parent_dir)
        print("===== Creating ZIP file")
        zipf = zipfile.ZipFile(os.path.join(
            parent_dir, f"Course – {cp_name} {position} Interview preparation.zip"), 'w', zipfile.ZIP_DEFLATED)
        zipdir(os.path.join(parent_dir, "Course"), zipf)
        zipf.close()
        rmtree(os.path.join(temp_dir, "sys_temp_dir"))
        print("===== Finished - Company: " +
              cp_name, "Position: " + position)
        print(" ")
        print(
            f"========== Create landing page for {cp_name} {position} Interview Preparation Online Course ")
        title = f"{cp_name} {position} Interview Preparation Online Course"
        print("===== Uploaded Course into Shopify ")
        _data = add_product(cp_name, position, industry,
                            price_course, title, shopify_copy_course)
        p_id = _data['product']['id']
        v_id = _data['product']['variants'][0]['id']

        update_requires_shiping(v_id)
        print("===== Updated variant is not phisical shipping ")

        (_, c_id) = check_collection(cp_name)
        if c_id is not None:
            print("===== Add Course to the collection ")
            add_product_to_collection(p_id, c_id)
        else:
            print("===== Create a collection for Course ")
            create_collection(cp_name, logo)
            print("===== Add Course to the collection ")
            add_product_to_collection(p_id, c_id)
        print(f"===== Uploading Course's image to shopify with product ")
        img_path_course = os.path.join(
            templates, "Images", "Course", f"{cp_name} {position}.jpg")
        b = upload_image(p_id, img_path_course)

        zip_path = os.path.join(
            parent_dir, f"Course – {cp_name} {position} Interview preparation.zip")
        print("===== Uploading Course's Zip file to Sendowl ")
        owl = send_owl(title, price, v_id, zip_path, pdf_stamping=False)
        print(f"========== Uploaded! \n")

        print(
            f"========== Create Book's landing page for {cp_name} {position} ")

        title = f"{cp_name} {position} Interview Preparation Study Guide"
        _data = add_product(cp_name, position, industry,
                            price, title, shopify_copy_book)
        p_id = _data['product']['id']
        v_id = _data['product']['variants'][0]['id']

        print(f"===== Added Book to the collection {cp_name}")
        print(f"===== Uploaded {title} into Shopify ")

        (_, c_id) = check_collection(cp_name)

        if c_id is not None:

            add_product_to_collection(p_id, c_id)
            print(f"===== Added Book to the collection {cp_name} ")
        else:
            create_collection(cp_name, logo)
            print(f"===== Created a collection of Book for {cp_name} ")
            add_product_to_collection(p_id, c_id)
            print(f"===== Added Book to the collection {cp_name} ")

        print(f"===== Uploading Book's image to shopify with product ")
        img_path_book = os.path.join(
            templates, "Images", "Study Guide", f"{cp_name} {position}.jpg")
        b = upload_image(p_id, img_path_book)
        pdf_path = os.path.join(
            parent_dir, "Study Guide", f'Study Guide–{cp_name} {position} Interview preparation.pdf')
        print("===== Uploading Book's Pdf file to Sendowl ")
        owl = send_owl(title, price, v_id, zip_path, pdf_stamping=True)
        print(f"========== Uploaded! \n")
